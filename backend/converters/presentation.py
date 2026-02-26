"""Conversor de apresentações PPTX → PDF / PNG via LibreOffice ou python-pptx."""
import subprocess
import shutil
from pathlib import Path


def convert(input_path: str, output_path: str, target_format: str) -> None:
    target_format = target_format.lower()

    if target_format == "pdf":
        _to_pdf(input_path, output_path)
    elif target_format == "png":
        _to_png(input_path, output_path)
    else:
        raise ValueError(f"PPTX não suporta saída: {target_format}")


def _to_pdf(input_path: str, output_path: str) -> None:
    """Converte PPTX → PDF. Tenta LibreOffice, depois fallback texto."""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if lo:
        import tempfile, os
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                [lo, "--headless", "--convert-to", "pdf",
                 "--outdir", tmpdir, input_path],
                capture_output=True, text=True
            )
            # LibreOffice gera o PDF com o mesmo nome base
            stem = Path(input_path).stem
            generated = Path(tmpdir) / f"{stem}.pdf"
            if generated.exists():
                generated.rename(output_path)
                return
    # Fallback: extrair texto e gerar PDF simples
    _pptx_text_to_pdf(input_path, output_path)


def _to_png(input_path: str, output_path: str) -> None:
    """Converte PPTX → PNG (slides combinados verticalmente)."""
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if lo:
        import tempfile
        with tempfile.TemporaryDirectory() as tmpdir:
            # Converter para PDF primeiro, depois PDF → PNG
            result = subprocess.run(
                [lo, "--headless", "--convert-to", "pdf",
                 "--outdir", tmpdir, input_path],
                capture_output=True, text=True
            )
            stem = Path(input_path).stem
            pdf_path = Path(tmpdir) / f"{stem}.pdf"
            if pdf_path.exists():
                from converters.document import _pdf_convert
                _pdf_convert(str(pdf_path), output_path, "png")
                return

    # Fallback: extrair texto, renderizar como PNG de texto
    _pptx_text_to_png(input_path, output_path)


def _pptx_text_to_pdf(input_path: str, output_path: str) -> None:
    """Extrai texto do PPTX e gera PDF simples."""
    from pptx import Presentation
    import weasyprint

    prs = Presentation(input_path)
    html_parts = [
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        "<style>body{font-family:sans-serif;margin:40px} "
        ".slide{page-break-after:always;border:1px solid #ccc;padding:20px;margin-bottom:20px}"
        "h2{color:#333}</style></head><body>"
    ]
    for i, slide in enumerate(prs.slides, 1):
        html_parts.append(f"<div class='slide'><h2>Slide {i}</h2>")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                html_parts.append(f"<p>{shape.text}</p>")
        html_parts.append("</div>")
    html_parts.append("</body></html>")
    weasyprint.HTML(string="\n".join(html_parts)).write_pdf(output_path)


def _pptx_text_to_png(input_path: str, output_path: str) -> None:
    """Extrai texto do PPTX e gera PNG simples."""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(input_path)
    slide_images = []
    W, H = 800, 600

    for i, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)
        draw.text((20, 20), f"Slide {i}", fill=(50, 50, 200))
        y = 60
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                for line in shape.text.splitlines():
                    draw.text((20, y), line[:90], fill=(30, 30, 30))
                    y += 22
                    if y > H - 30:
                        break
        slide_images.append(img)

    if not slide_images:
        raise RuntimeError("Apresentação sem slides")

    total_h = sum(img.height for img in slide_images)
    combined = Image.new("RGB", (W, total_h), "white")
    y = 0
    for img in slide_images:
        combined.paste(img, (0, y))
        y += img.height
    combined.save(output_path, format="PNG")
